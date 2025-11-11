#!/usr/bin/env python3
"""
Generador de presentación PowerPoint
Basado en el diseño del PDF Impresion3D_Vurokrazia_v2.pptx.pdf
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Colores del diseño (extraídos del PDF)
COLOR_FONDO = RGBColor(45, 55, 72)  # Azul marino oscuro #2d3748
COLOR_TITULO_PRINCIPAL = RGBColor(159, 122, 234)  # Morado #9f7aea
COLOR_TITULO_SECUNDARIO = RGBColor(147, 197, 253)  # Azul claro
COLOR_TEXTO = RGBColor(255, 255, 255)  # Blanco
COLOR_ACENTO = RGBColor(159, 122, 234)  # Morado para líneas
COLOR_VERDE = RGBColor(79, 209, 197)  # Verde cyan #4fd1c5
COLOR_ROJO = RGBColor(245, 101, 101)  # Rojo salmon #f56565

# Colores para cajas
COLOR_CAJA_MITO_FONDO = RGBColor(74, 47, 58)  # Burgundy oscuro
COLOR_CAJA_MITO_BORDE = RGBColor(245, 101, 101)  # Rojo
COLOR_CAJA_REALIDAD_FONDO = RGBColor(31, 79, 82)  # Teal oscuro
COLOR_CAJA_REALIDAD_BORDE = RGBColor(79, 209, 197)  # Cyan


def crear_presentacion():
    """Crea la presentación completa"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Portada
    crear_slide_portada(prs)

    # SLIDE 2: Sobre Mí
    crear_slide_sobre_mi(prs)

    # SLIDE 3: ¿Por Qué Ahora?
    crear_slide_por_que_ahora(prs)

    # SLIDE 4: Separador Mitos vs Realidad
    crear_slide_separador_mitos(prs)

    # SLIDES 5-20: Los 8 mitos (2 slides por mito)
    crear_slides_mitos(prs)

    # SLIDES 21-25: Conceptos Básicos FDM
    crear_slides_conceptos_fdm(prs)

    # SLIDES 26-30: Los Materiales
    crear_slides_materiales(prs)

    # SLIDES 31-35: Repositorios
    crear_slides_repositorios(prs)

    # SLIDES 36-40: Eligiendo Impresora
    crear_slides_eligiendo_impresora(prs)

    # SLIDES 41-44: Software
    crear_slides_software(prs)

    # SLIDES 45-50: Primeros Pasos
    crear_slides_primeros_pasos(prs)

    # SLIDE 51: Gracias
    crear_slide_gracias(prs)

    # Guardar presentación
    prs.save('presentacion/Impresion3D_Vurokrazia_Completa.pptx')
    print("✅ Presentación generada: presentacion/Impresion3D_Vurokrazia_Completa.pptx")


def crear_fondo(slide):
    """Agrega fondo oscuro al slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_FONDO


def agregar_linea_acento(slide, top):
    """Agrega línea morada horizontal de acento"""
    line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(4.5), Inches(top),
        Inches(1), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACENTO
    line.line.fill.background()


def crear_slide_portada(prs):
    """Slide 1: Portada"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    crear_fondo(slide)

    # Línea de acento superior
    agregar_linea_acento(slide, 1.5)

    # Título principal
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    tf = titulo.text_frame
    tf.text = "Entrar al Mundo de la\nImpresión 3D"
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(60)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO
        paragraph.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
    tf = subtitulo.text_frame
    tf.text = "Sin Fallar en el Intento"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Autor
    autor = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = autor.text_frame
    tf.text = "Jesús Martínez • Vurokrazia3D"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


def crear_slide_sobre_mi(prs):
    """Slide 2: Sobre Mí"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "Sobre Mí"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL

    # Caja con borde izquierdo - Jesús Martínez
    y_pos = 1.5
    borde = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(0.08), Inches(0.8))
    borde.fill.solid()
    borde.fill.fore_color.rgb = COLOR_TITULO_PRINCIPAL
    borde.line.fill.background()

    contenido = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.5), Inches(0.8))
    tf = contenido.text_frame
    tf.text = "Jesús Martínez\nProgramador con más de 7 años de experiencia"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_TEXTO
    tf.paragraphs[1].font.size = Pt(20)
    tf.paragraphs[1].font.color.rgb = COLOR_TEXTO

    # Vurokrazia3D
    y_pos = 2.8
    titulo2 = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5))
    tf = titulo2.text_frame
    tf.text = "Vurokrazia3D"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE

    desc = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.5), Inches(9), Inches(0.8))
    tf = desc.text_frame
    tf.text = "Proyecto personal donde combino programación con impresión 3D y electrónica Arduino"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXTO

    # Enfoque
    y_pos = 4.3
    titulo3 = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5))
    tf = titulo3.text_frame
    tf.text = "Enfoque"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL

    bullets = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.6), Inches(9), Inches(2))
    tf = bullets.text_frame
    items = [
        "Diseño y creación de piezas",
        "Acabado profesional",
        "Integración con sistemas electrónicos",
        "Materializar ideas técnicas y creativas"
    ]
    for i, item in enumerate(items):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO


def crear_slide_por_que_ahora(prs):
    """Slide 3: ¿Por Qué Ahora es el Momento?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "¿Por Qué Ahora es el Momento?"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.5))
    tf = subtitulo.text_frame
    tf.text = "La impresión 3D ya no es solo para empresas o expertos"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_TEXTO

    # Tres cajas
    cajas = [
        {"titulo": "$3,500 MXN", "sub": "($200 USD)", "desc": "Impresoras desde"},
        {"titulo": "Millones", "sub": "", "desc": "de modelos gratuitos"},
        {"titulo": "2025", "sub": "", "desc": "Tecnología madura"}
    ]

    x_start = 0.5
    y_pos = 2.5
    box_width = 2.8
    spacing = 0.2

    for i, caja in enumerate(cajas):
        x = x_start + i * (box_width + spacing)

        # Línea superior
        line = slide.shapes.add_shape(1, Inches(x), Inches(y_pos), Inches(box_width), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_TITULO_PRINCIPAL
        line.line.fill.background()

        # Contenido
        content = slide.shapes.add_textbox(Inches(x), Inches(y_pos + 0.3), Inches(box_width), Inches(2))
        tf = content.text_frame

        # Título grande
        tf.text = caja["titulo"]
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_VERDE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Subtítulo
        if caja["sub"]:
            tf.add_paragraph()
            tf.paragraphs[1].text = caja["sub"]
            tf.paragraphs[1].font.size = Pt(16)
            tf.paragraphs[1].font.color.rgb = COLOR_TEXTO
            tf.paragraphs[1].alignment = PP_ALIGN.CENTER

        # Descripción
        tf.add_paragraph()
        last_p = tf.paragraphs[-1]
        last_p.text = caja["desc"]
        last_p.font.size = Pt(18)
        last_p.font.color.rgb = COLOR_TEXTO
        last_p.alignment = PP_ALIGN.CENTER

    # Texto inferior
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = footer.text_frame
    tf.text = "Software gratuito • Comunidad activa • Fácil de comenzar"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


def crear_slide_separador_mitos(prs):
    """Slide 4: Separador Mitos vs Realidad"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Línea de acento
    agregar_linea_acento(slide, 2)

    # Título principal
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = titulo.text_frame
    tf.text = "Mitos vs Realidad"
    p = tf.paragraphs[0]
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "Desmintiendo las creencias más comunes"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


def crear_slides_mitos(prs):
    """Crea los slides de los 8 mitos"""
    mitos = [
        {
            "num": 1,
            "mito": "Necesitas saber diseñar en 3D para empezar",
            "realidad": "Millones de modelos gratuitos listos para descargar",
            "detalle": "Thingiverse, Printables, MyMiniFactory tienen modelos listos. Puedes empezar a imprimir el mismo día sin diseñar nada."
        },
        {
            "num": 2,
            "mito": "Todas las impresoras son iguales, solo cambia el precio",
            "realidad": "El precio determina QUÉ puedes imprimir",
            "detalle": "$5,500 MXN ($300 USD): solo PLA básico | $9,000-18,000 MXN ($500-1000 USD): ABS, PETG | $18,000+ MXN ($1000+ USD): Nylon, fibra de carbono"
        },
        {
            "num": 3,
            "mito": "La impresión 3D es cara y solo para empresas",
            "realidad": "Puedes empezar con menos de $5,500 MXN ($300 USD)",
            "detalle": "Impresoras desde $3,500 MXN ($200 USD) • Filamento PLA desde $320 MXN ($18 USD) el kg • Software gratuito"
        },
        {
            "num": 4,
            "mito": "Es muy complicado, hay que ser ingeniero",
            "realidad": "Las impresoras 2025 son plug & play",
            "detalle": "Nivelación automática • Perfiles preconfigurados • Software gratuito e intuitivo (Cura) • Tutoriales en español • Comunidad activa"
        },
        {
            "num": 5,
            "mito": "Los plásticos son todos iguales",
            "realidad": "Cada material tiene propiedades únicas",
            "detalle": "PLA: fácil pero frágil • ABS: resistente pero difícil • PETG: equilibrado • TPU: flexible • Nylon: industrial"
        },
        {
            "num": 6,
            "mito": "Necesitas una oficina o taller especial",
            "realidad": "Con PLA puedes imprimir en tu escritorio",
            "detalle": "PLA no huele ni es tóxico • Solo ABS requiere ventilación • Ruido moderado • Tamaño compacto"
        },
        {
            "num": 7,
            "mito": "Las piezas son débiles y se rompen fácilmente",
            "realidad": "Depende del material y configuración",
            "detalle": "ABS y PETG son muy resistentes • Se usan en autos y equipos industriales • Configuración del relleno importa • Piezas funcionales reales"
        },
        {
            "num": 8,
            "mito": "Toma años dominar la impresión 3D",
            "realidad": "Primera impresión exitosa en horas, no meses",
            "detalle": "Máquinas modernas auto-calibran • Perfiles listos para usar • Tutoriales abundantes • Comunidad ayuda activamente"
        }
    ]

    for mito_data in mitos:
        # Slide solo con MITO
        crear_slide_mito_solo(prs, mito_data)
        # Slide con MITO + REALIDAD
        crear_slide_mito_completo(prs, mito_data)


def crear_slide_mito_solo(prs, mito_data):
    """Crea slide solo con el MITO"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Número del mito (pequeño arriba)
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    tf = num_box.text_frame
    tf.text = f"MITO #{mito_data['num']}"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXTO

    # Caja grande con el mito
    crear_caja_mito(slide, mito_data['mito'], 2.5)


def crear_slide_mito_completo(prs, mito_data):
    """Crea slide con MITO + REALIDAD"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Número del mito
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    tf = num_box.text_frame
    tf.text = f"MITO #{mito_data['num']}"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXTO

    # Caja MITO (más pequeña)
    crear_caja_mito(slide, mito_data['mito'], 1.2, altura=1.2)

    # Caja REALIDAD
    crear_caja_realidad(slide, mito_data['realidad'], mito_data['detalle'], 3.0)


def crear_caja_mito(slide, texto, y_pos, altura=2.5):
    """Crea caja estilo MITO con borde izquierdo"""
    # Fondo
    fondo = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(8.5), Inches(altura))
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = COLOR_CAJA_MITO_FONDO
    fondo.line.fill.background()

    # Borde izquierdo
    borde = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(0.08), Inches(altura))
    borde.fill.solid()
    borde.fill.fore_color.rgb = COLOR_CAJA_MITO_BORDE
    borde.line.fill.background()

    # Texto "❌ MITO"
    label = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.15), Inches(7.5), Inches(0.4))
    tf = label.text_frame
    tf.text = "❌ MITO"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_ROJO

    # Texto principal
    content = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.5), Inches(7.5), Inches(altura - 0.6))
    tf = content.text_frame
    tf.text = texto
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROJO


def crear_caja_realidad(slide, titulo, detalle, y_pos):
    """Crea caja estilo REALIDAD con borde izquierdo"""
    altura = 2.5

    # Fondo
    fondo = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(8.5), Inches(altura))
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = COLOR_CAJA_REALIDAD_FONDO
    fondo.line.fill.background()

    # Borde izquierdo
    borde = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(0.08), Inches(altura))
    borde.fill.solid()
    borde.fill.fore_color.rgb = COLOR_CAJA_REALIDAD_BORDE
    borde.line.fill.background()

    # Texto "✅ REALIDAD"
    label = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.15), Inches(7.5), Inches(0.4))
    tf = label.text_frame
    tf.text = "✅ REALIDAD"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_VERDE

    # Título principal
    titulo_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.5), Inches(7.5), Inches(0.8))
    tf = titulo_box.text_frame
    tf.text = titulo
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE

    # Detalle
    detalle_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 1.4), Inches(7.5), Inches(0.9))
    tf = detalle_box.text_frame
    tf.text = detalle
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXTO


def crear_slides_conceptos_fdm(prs):
    """Crea los slides de Conceptos Básicos FDM"""
    # Slide 21: Separador Conceptos Básicos FDM
    crear_slide_separador_conceptos(prs)

    # Slide 22: La impresora derrite plástico
    crear_slide_derrite_plastico(prs)

    # Slide 23: Temperaturas típicas
    crear_slide_temperaturas(prs)

    # Slide 24: Movimiento en 3 ejes
    crear_slide_ejes(prs)

    # Slide 25: La primera capa es CRUCIAL
    crear_slide_primera_capa(prs)


def crear_slide_separador_conceptos(prs):
    """Slide 21: Separador Conceptos Básicos FDM"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título principal con emoji
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = titulo.text_frame
    tf.text = "🧩 Conceptos Básicos FDM"
    p = tf.paragraphs[0]
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "¿Cómo funciona?"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


def crear_slide_derrite_plastico(prs):
    """Slide 22: La impresora derrite plástico"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Emoji grande
    emoji = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    tf = emoji.text_frame
    tf.text = "🔥 La impresora derrite plástico"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo grande
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.8))
    tf = subtitulo.text_frame
    tf.text = "Capa por capa"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Descripción
    desc = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(1))
    tf = desc.text_frame
    tf.text = "El filamento se calienta y deposita línea por línea"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


def crear_slide_temperaturas(prs):
    """Slide 23: Temperaturas típicas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🌡️ Temperaturas típicas"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Lista de temperaturas
    materiales = [
        ("PLA:", "200°C boquilla / 60°C cama"),
        ("PETG:", "240°C boquilla / 80°C cama"),
        ("ABS:", "250°C boquilla / 100°C cama")
    ]

    y_start = 2.5
    spacing = 1

    for i, (material, temps) in enumerate(materiales):
        y_pos = y_start + i * spacing

        # Contenedor para cada material
        content = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.8))
        tf = content.text_frame

        # Material en bold
        tf.text = material
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_VERDE

        # Temperaturas
        tf.add_paragraph()
        tf.paragraphs[1].text = temps
        tf.paragraphs[1].font.size = Pt(24)
        tf.paragraphs[1].font.color.rgb = COLOR_TEXTO


def crear_slide_ejes(prs):
    """Slide 24: Movimiento en 3 ejes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "⚙️ Movimiento en 3 ejes"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Lista de ejes
    ejes = [
        ("Eje X:", "izquierda ↔ derecha"),
        ("Eje Y:", "adelante ↔ atrás"),
        ("Eje Z:", "arriba ↕ abajo")
    ]

    y_start = 2.2
    spacing = 0.9

    for i, (eje, movimiento) in enumerate(ejes):
        y_pos = y_start + i * spacing

        # Bullet point
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.7))
        tf = bullet.text_frame

        tf.text = f"• {eje} {movimiento}"
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXTO

    # Subtítulo final
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "Precisión milimétrica"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER


def crear_slide_primera_capa(prs):
    """Slide 25: La primera capa es CRUCIAL"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título con emoji
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = titulo.text_frame
    tf.text = "✅ La primera capa es CRUCIAL"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Mensaje importante
    mensaje = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
    tf = mensaje.text_frame
    tf.text = "Si falla la primera capa, falla todo"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


def crear_slides_materiales(prs):
    """Crea los slides de Los Materiales"""
    # Slide 26: Separador Los Materiales
    crear_slide_separador_materiales(prs)

    # Slide 27: PLA - El principiante
    crear_slide_material_pla(prs)

    # Slide 28: PETG - El equilibrado
    crear_slide_material_petg(prs)

    # Slide 29: TPU - El flexible
    crear_slide_material_tpu(prs)

    # Slide 30: ABS - El resistente
    crear_slide_material_abs(prs)


def crear_slide_separador_materiales(prs):
    """Slide 26: Separador Los Materiales"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título principal
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = titulo.text_frame
    tf.text = "🎨 Los Materiales"
    p = tf.paragraphs[0]
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "Comparativa rápida"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


def crear_slide_material_pla(prs):
    """Slide 27: PLA - El principiante"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🟢 PLA - El principiante"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Ventajas
    ventajas = [
        "✅ MÁS FÁCIL de imprimir",
        "✅ No huele, no tóxico",
        "✅ Biodegradable (maíz)"
    ]

    y_pos = 2.0
    for ventaja in ventajas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = ventaja
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.6

    # Desventajas
    desventajas = [
        "❌ Frágil ante golpes",
        "❌ Se ablanda con calor (+60°C)"
    ]

    for desventaja in desventajas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = desventaja
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.6

    # Precio
    precio = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.6))
    tf = precio.text_frame
    tf.text = "💰 Desde $320 MXN/kg"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER


def crear_slide_material_petg(prs):
    """Slide 28: PETG - El equilibrado"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🔷 PETG - El equilibrado"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_SECUNDARIO
    p.alignment = PP_ALIGN.CENTER

    # Características
    caracteristicas = [
        "✅ Resistente a impactos",
        "✅ Resistente al agua",
        "✅ Buena resistencia térmica",
        "⚠️ Genera hilos (stringing)"
    ]

    y_pos = 2.2
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.7

    # Recomendación
    recomendacion = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(0.8))
    tf = recomendacion.text_frame
    tf.text = "🏆 Recomendado para piezas funcionales"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER


def crear_slide_material_tpu(prs):
    """Slide 29: TPU - El flexible"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🟣 TPU - El flexible"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Características
    caracteristicas = [
        "✅ Súper flexible (tipo goma)",
        "✅ Absorbe impactos",
        "✅ No se rompe",
        "❌ Difícil de imprimir rápido"
    ]

    y_pos = 2.2
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.7

    # Ideal para
    ideal = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(0.8))
    tf = ideal.text_frame
    tf.text = "🎯 Ideal para: llantas, sellos, fundas"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER


def crear_slide_material_abs(prs):
    """Slide 30: ABS - El resistente"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🟠 ABS - El resistente"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(251, 146, 60)  # Naranja
    p.alignment = PP_ALIGN.CENTER

    # Ventajas
    ventajas = [
        "✅ Muy resistente",
        "✅ Soporta calor",
        "✅ Se usa en LEGO"
    ]

    y_pos = 2.0
    for ventaja in ventajas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = ventaja
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.6

    # Desventajas
    desventajas = [
        "❌ Difícil (warping)",
        "❌ Huele mal",
        "❌ Requiere ventilación"
    ]

    for desventaja in desventajas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = desventaja
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.6

    # Advertencia
    advertencia = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.6))
    tf = advertencia.text_frame
    tf.text = "⚠️ NO recomendado para principiantes"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROJO
    p.alignment = PP_ALIGN.CENTER


def crear_slides_repositorios(prs):
    """Crea los slides de Repositorios"""
    # Slide 31: Separador Repositorios
    crear_slide_separador_repositorios(prs)

    # Slide 32: Thingiverse
    crear_slide_thingiverse(prs)

    # Slide 33: Printables
    crear_slide_printables(prs)

    # Slide 34: MyMiniFactory
    crear_slide_mymini(prs)

    # Slide 35: Puedes empezar HOY
    crear_slide_empezar_hoy(prs)


def crear_slide_separador_repositorios(prs):
    """Slide 31: Separador Repositorios"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título principal
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.2))
    tf = titulo.text_frame
    tf.text = "🏛️ ¿Dónde Conseguir Modelos?"
    p = tf.paragraphs[0]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER


def crear_slide_thingiverse(prs):
    """Slide 32: Thingiverse"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🟦 Thingiverse"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)  # Azul
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "El más grande y antiguo"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Características
    caracteristicas = [
        "• Millones de modelos",
        "• Comunidad masiva",
        "• Todo tipo de categorías"
    ]

    y_pos = 3.8
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(5), Inches(0.5))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(26)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.7


def crear_slide_printables(prs):
    """Slide 33: Printables"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🟩 Printables"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "Por Prusa - Moderno y organizado"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Características
    caracteristicas = [
        "• Modelos de calidad",
        "• Sistema de puntuación",
        "• Concursos regulares"
    ]

    y_pos = 3.8
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(5), Inches(0.5))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(26)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.7


def crear_slide_mymini(prs):
    """Slide 34: MyMiniFactory"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🟧 MyMiniFactory"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(251, 146, 60)  # Naranja
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "Enfocado en calidad"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Características
    caracteristicas = [
        "• Todos los modelos son probados",
        "• Excelente para miniaturas"
    ]

    y_pos = 3.8
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(5), Inches(0.5))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(26)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.7


def crear_slide_empezar_hoy(prs):
    """Slide 35: Puedes empezar HOY"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Mensaje principal
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = titulo.text_frame
    tf.text = "✨ Puedes empezar a imprimir HOY"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    tf = subtitulo.text_frame
    tf.text = "Sin diseñar nada"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


def crear_slides_eligiendo_impresora(prs):
    """Crea los slides de Eligiendo Impresora"""
    # Slide 36: Separador
    crear_slide_separador_impresora(prs)

    # Slide 37: Gama Entrada
    crear_slide_gama_entrada(prs)

    # Slide 38: Gama Media
    crear_slide_gama_media(prs)

    # Slide 39: Gama Alta
    crear_slide_gama_alta(prs)

    # Slide 40: Recomendación
    crear_slide_recomendacion_impresora(prs)


def crear_slide_separador_impresora(prs):
    """Slide 36: Separador Eligiendo Impresora"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título principal
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.2))
    tf = titulo.text_frame
    tf.text = "📦 Eligiendo tu Primera Impresora"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER


def crear_slide_gama_entrada(prs):
    """Slide 37: Gama Entrada"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.6))
    tf = titulo.text_frame
    tf.text = "🟢 Gama Entrada"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Precio
    precio = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    tf = precio.text_frame
    tf.text = "$3,500 - $5,500 MXN ($200-$300 USD)"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Lo que obtienes
    label1 = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(7), Inches(0.4))
    tf = label1.text_frame
    tf.text = "Lo que obtienes:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE

    caracteristicas = [
        "• Solo PLA",
        "• Nivelación manual/semiautomática",
        "• Velocidades moderadas"
    ]

    y_pos = 2.8
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.5

    # Ideal para
    label2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(0.4))
    tf = label2.text_frame
    tf.text = "Ideal para:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE

    ideal = [
        "• Principiantes absolutos",
        "• Figuras decorativas"
    ]

    y_pos = 5.3
    for item in ideal:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = bullet.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.5


def crear_slide_gama_media(prs):
    """Slide 38: Gama Media"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.6))
    tf = titulo.text_frame
    tf.text = "🟡 Gama Media"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(234, 179, 8)  # Amarillo
    p.alignment = PP_ALIGN.CENTER

    # Precio
    precio = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    tf = precio.text_frame
    tf.text = "$9,000 - $18,000 MXN ($500-$1000 USD)"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Lo que obtienes
    label1 = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(7), Inches(0.4))
    tf = label1.text_frame
    tf.text = "Lo que obtienes:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(234, 179, 8)

    caracteristicas = [
        "• PLA, PETG, ABS",
        "• Nivelación automática",
        "• Velocidades más rápidas",
        "• Mejor calidad"
    ]

    y_pos = 2.8
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.5

    # Ideal para
    label2 = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(0.4))
    tf = label2.text_frame
    tf.text = "Ideal para:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(234, 179, 8)

    ideal = [
        "• Piezas funcionales",
        "• Prototipos de calidad"
    ]

    y_pos = 5.5
    for item in ideal:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = bullet.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.5


def crear_slide_gama_alta(prs):
    """Slide 39: Gama Alta"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.6))
    tf = titulo.text_frame
    tf.text = "🔴 Gama Alta"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROJO
    p.alignment = PP_ALIGN.CENTER

    # Precio
    precio = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    tf = precio.text_frame
    tf.text = "$18,000+ MXN ($1000+ USD)"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Lo que obtienes
    label1 = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(7), Inches(0.4))
    tf = label1.text_frame
    tf.text = "Lo que obtienes:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROJO

    caracteristicas = [
        "• Todos los materiales (Nylon, fibra de carbono)",
        "• Cámara cerrada",
        "• Velocidades súper altas (500mm/s+)",
        "• Máxima precisión"
    ]

    y_pos = 2.8
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.5

    # Ideal para
    label2 = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(0.4))
    tf = label2.text_frame
    tf.text = "Ideal para:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROJO

    ideal = [
        "• Profesionales",
        "• Producción"
    ]

    y_pos = 5.5
    for item in ideal:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = bullet.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.5


def crear_slide_recomendacion_impresora(prs):
    """Slide 40: Recomendación"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🎯 Recomendación"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.5))
    tf = subtitulo.text_frame
    tf.text = "Para tu primera impresora"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Recomendaciones
    recomendaciones = [
        "• Presupuesto mínimo: $5,500 MXN",
        "• Busca: Nivelación automática",
        "• Empieza con: PLA",
        "• Investiga: YouTube reviews en español"
    ]

    y_pos = 3.0
    for rec in recomendaciones:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = rec
        p = tf.paragraphs[0]
        p.font.size = Pt(26)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.8


def crear_slides_software(prs):
    """Crea los slides de Software"""
    # Slide 41: Separador Software
    crear_slide_separador_software(prs)

    # Slide 42: Cura
    crear_slide_cura(prs)

    # Slide 43: Flujo de Trabajo
    crear_slide_flujo_trabajo(prs)

    # Slide 44: Configuración Recomendada
    crear_slide_config_recomendada(prs)


def crear_slide_separador_software(prs):
    """Slide 41: Separador Software"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título principal
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.2))
    tf = titulo.text_frame
    tf.text = "💻 Software Básico"
    p = tf.paragraphs[0]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER


def crear_slide_cura(prs):
    """Slide 42: Cura (Ultimaker)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.8))
    tf = titulo.text_frame
    tf.text = "🌐 Cura (Ultimaker)"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "El más popular - 100% GRATUITO"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Características
    caracteristicas = [
        "✅ Interfaz intuitiva",
        "✅ Perfiles pre-configurados",
        "✅ Gran comunidad",
        "✅ Windows, Mac, Linux"
    ]

    y_pos = 3.5
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(5), Inches(0.5))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(26)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.7


def crear_slide_flujo_trabajo(prs):
    """Slide 43: Flujo de Trabajo Simple"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.7))
    tf = titulo.text_frame
    tf.text = "📊 Flujo de Trabajo Simple"
    p = tf.paragraphs[0]
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Pasos
    pasos = [
        ("1️⃣", "Descargar modelo (.STL)", "Thingiverse / Printables"),
        ("2️⃣", "Abrir en Cura", "Arrastra el archivo"),
        ("3️⃣", "Configurar", "Material, calidad, relleno"),
        ("4️⃣", "Generar G-code", '"Slice" / "Laminar"'),
        ("5️⃣", "Copiar a SD", "O enviar por WiFi"),
        ("6️⃣", "¡Imprimir!", "")
    ]

    y_pos = 1.6
    for emoji, paso, detalle in pasos:
        # Número/Emoji
        num = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(0.8), Inches(0.4))
        tf = num.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.font.size = Pt(22)

        # Paso
        paso_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(3), Inches(0.4))
        tf = paso_box.text_frame
        tf.text = paso
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXTO

        # Detalle
        if detalle:
            detalle_box = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(3), Inches(0.4))
            tf = detalle_box.text_frame
            tf.text = detalle
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(156, 163, 175)  # Gris

        y_pos += 0.75


def crear_slide_config_recomendada(prs):
    """Slide 44: Configuración Recomendada"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(8), Inches(0.7))
    tf = titulo.text_frame
    tf.text = "⚙️ Configuración Recomendada"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    tf = subtitulo.text_frame
    tf.text = "Para tu primera impresión"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER

    # Configuración (simulando código)
    config_items = [
        "Material: PLA",
        "Temperatura boquilla: 200°C",
        "Temperatura cama: 60°C",
        "Velocidad: 50mm/s",
        "Altura de capa: 0.2mm",
        "Relleno: 20%",
        "Soportes: NO (elige modelo sin voladizos)"
    ]

    # Caja de fondo para simular código
    fondo_codigo = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5), Inches(2.3),
        Inches(7), Inches(3.8)
    )
    fondo_codigo.fill.solid()
    fondo_codigo.fill.fore_color.rgb = RGBColor(30, 41, 59)  # Azul oscuro
    fondo_codigo.line.color.rgb = COLOR_TITULO_PRINCIPAL
    fondo_codigo.line.width = Pt(2)

    y_pos = 2.5
    for config in config_items:
        config_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = config_box.text_frame
        tf.text = config
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.name = "Courier New"
        p.font.color.rgb = COLOR_VERDE
        y_pos += 0.5


def crear_slides_primeros_pasos(prs):
    """Crea los slides de Primeros Pasos"""
    # Slide 45: Separador
    crear_slide_separador_primeros_pasos(prs)

    # Slide 46: Paso 1 - Calibración
    crear_slide_paso1_calibracion(prs)

    # Slide 47: Paso 2 - Tu Primera Pieza
    crear_slide_paso2_primera_pieza(prs)

    # Slide 48: Paso 3 - Prepárate
    crear_slide_paso3_preparate(prs)

    # Slide 49: Paso 4 - Observa
    crear_slide_paso4_observa(prs)

    # Slide 50: Paso 5 - Paciencia
    crear_slide_paso5_paciencia(prs)


def crear_slide_separador_primeros_pasos(prs):
    """Slide 45: Separador Primeros Pasos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título principal
    titulo = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.2))
    tf = titulo.text_frame
    tf.text = "🚀 Tu Primera Impresión"
    p = tf.paragraphs[0]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER


def crear_slide_paso1_calibracion(prs):
    """Slide 46: Paso 1 - Calibración"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(8), Inches(0.7))
    tf = titulo.text_frame
    tf.text = "🎯 Paso 1: Calibración"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    tf = subtitulo.text_frame
    tf.text = "Nivelación de la cama"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Automática
    auto_label = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(7), Inches(0.4))
    tf = auto_label.text_frame
    tf.text = "Automática:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE

    auto_items = ["• Ejecuta desde el menú", "• ¡Listo!"]
    y_pos = 2.8
    for item in auto_items:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = bullet.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.5

    # Manual
    manual_label = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(0.4))
    tf = manual_label.text_frame
    tf.text = "Manual:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(234, 179, 8)  # Amarillo

    manual_items = [
        "• Papel entre boquilla y cama",
        "• Ajusta hasta que roce ligeramente",
        "• Repite en las 4 esquinas"
    ]
    y_pos = 4.5
    for item in manual_items:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.4))
        tf = bullet.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.5

    # Mensaje importante
    mensaje = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.6))
    tf = mensaje.text_frame
    tf.text = "✅ La primera capa es la BASE de todo"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER


def crear_slide_paso2_primera_pieza(prs):
    """Slide 47: Paso 2 - Tu Primera Pieza"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.7))
    tf = titulo.text_frame
    tf.text = "📝 Paso 2: Tu Primera Pieza"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "Elige algo SIMPLE"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Características
    label = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(0.4))
    tf = label.text_frame
    tf.text = "Características:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO

    caracteristicas = [
        "• Sin voladizos complicados",
        "• Tamaño pequeño (menos de 2 horas)",
        "• Ejemplos: cubo de calibración, llavero"
    ]
    y_pos = 3.5
    for caract in caracteristicas:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = caract
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.7


def crear_slide_paso3_preparate(prs):
    """Slide 48: Paso 3 - Prepárate"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.7))
    tf = titulo.text_frame
    tf.text = "🏁 Paso 3: Prepárate"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Lista de preparación
    items = [
        "✅ Filamento PLA cargado",
        "✅ Cama limpia (alcohol isopropílico)",
        "✅ Tarjeta SD lista",
        "✅ Configura en Cura (perfil PLA básico)"
    ]

    y_pos = 2.8
    for item in items:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(26)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.9


def crear_slide_paso4_observa(prs):
    """Slide 49: Paso 4 - Observa"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.7))
    tf = titulo.text_frame
    tf.text = "👀 Paso 4: Observa"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(0.6))
    tf = subtitulo.text_frame
    tf.text = "Los primeros 5 minutos son críticos"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROJO
    p.alignment = PP_ALIGN.CENTER

    # Verifica
    label = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(0.4))
    tf = label.text_frame
    tf.text = "Verifica:"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE

    items = [
        "• Primera capa se adhiere bien",
        "• Líneas uniformes",
        "• Sin levantamientos"
    ]
    y_pos = 3.0
    for item in items:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = bullet.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.7

    # Mensaje
    mensaje = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    tf = mensaje.text_frame
    tf.text = "Si hay problemas: PAUSA y ajusta"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROJO
    p.alignment = PP_ALIGN.CENTER


def crear_slide_paso5_paciencia(prs):
    """Slide 50: Paso 5 - Paciencia"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Título
    titulo = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.7))
    tf = titulo.text_frame
    tf.text = "⏳ Paso 5: Paciencia"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.5))
    tf = subtitulo.text_frame
    tf.text = "Deja que termine"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER

    # Consejos
    items = [
        "• No abras puertas (corrientes de aire)",
        "• No muevas la impresora",
        "• Espera que la cama se enfríe antes de retirar"
    ]

    y_pos = 3.2
    for item in items:
        bullet = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.6))
        tf = bullet.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXTO
        y_pos += 0.9


def crear_slide_gracias(prs):
    """Slide 51: Gracias"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_fondo(slide)

    # Línea de acento
    agregar_linea_acento(slide, 1.2)

    # Título "¡Gracias!"
    titulo = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1))
    tf = titulo.text_frame
    tf.text = "¡Gracias!"
    p = tf.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITULO_PRINCIPAL
    p.alignment = PP_ALIGN.CENTER

    # Borde izquierdo para info
    borde = slide.shapes.add_shape(1, Inches(3.5), Inches(3.2), Inches(0.06), Inches(1.5))
    borde.fill.solid()
    borde.fill.fore_color.rgb = COLOR_TITULO_PRINCIPAL
    borde.line.fill.background()

    # Nombre
    nombre = slide.shapes.add_textbox(Inches(3.7), Inches(3.2), Inches(3), Inches(0.5))
    tf = nombre.text_frame
    tf.text = "Jesús Martínez"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO

    # Usuario
    usuario = slide.shapes.add_textbox(Inches(3.7), Inches(3.7), Inches(3), Inches(0.5))
    tf = usuario.text_frame
    tf.text = "@vurokrazia"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_VERDE

    # ¿Preguntas?
    preguntas = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.6))
    tf = preguntas.text_frame
    tf.text = "¿Preguntas?"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = COLOR_TEXTO
    p.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    crear_presentacion()
